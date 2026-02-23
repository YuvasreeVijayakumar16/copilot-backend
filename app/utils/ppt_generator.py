
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime
import os
import pandas as pd
from docx import Document
from openai import OpenAI
import re
import matplotlib.pyplot as plt

from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

import time
from pptx.enum.text import PP_ALIGN
import numpy as np
import textwrap

# === ✅ Initialize OpenAI
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# === ✅ Generate GPT Insights and Recommendations
def generate_insights(df: pd.DataFrame):
    prompt = f"Analyze this data and provide 3–5 business insights and 2–3 strong recommendations:\n\n{df.head(10).to_string(index=False)}"
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst generating insights for supply chain decision makers."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.2,
    )
    output = response.choices[0].message.content.strip()

    # Cleanup: remove markdown styling and numbered bullets
    output = re.sub(r"\*\*(.*?)\*\*", r"\1", output)
    output = re.sub(r"\*(.*?)\*", r"\1", output)
    output = re.sub(r"^\s*\d+\.\s*", "", output, flags=re.MULTILINE)

    parts = re.split(r"\n{2,}", output)
    insights, recs = "", ""
    for part in parts:
        if "recommend" in part.lower():
            recs = part.strip()
        else:
            insights = part.strip()
    return insights, recs

# === ✅ Create simple bar chart image
def create_bar_chart(df: pd.DataFrame, column: str = None, output_path: str = "generated_files/bar_chart.png") -> str:
    if not column:
        for col in df.columns:
            if df[col].dtype == "object" or df[col].nunique() < 20:
                column = col
                break
    if not column:
        return None

    chart_data = df[column].value_counts().head(5)
    plt.figure(figsize=(6, 4))
    chart_data.plot(kind="bar", color="skyblue")
    plt.title(f"Top 5 {column}")
    plt.xlabel(column)
    plt.ylabel("Count")
    plt.tight_layout()

    os.makedirs("generated_files", exist_ok=True)
    plt.savefig(output_path)
    plt.close()
    return output_path

def generate_unique_blob_name(prefix="PPT", actions_list=None, extension="pptx"):
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    unique_id = uuid.uuid4().hex
    if actions_list:
        snippet = "_".join([''.join(e for e in act if e.isalnum())[:20] for act in actions_list[:2]])
    else:
        snippet = "Output"
    return f"{prefix}_{snippet}_{timestamp}_{unique_id}.{extension}"

def extract_actions(message):
    parts = [p.strip() for p in re.split(r"\band\b|,|\s{2,}", message) if p.strip()]
    return parts

def select_descriptive_action(actions):
    if not actions:
        return ""
    sorted_actions = sorted(actions, key=lambda x: len(x), reverse=True)
    for act in sorted_actions:
        if len(act) > 10:
            return act
    return actions[0]

# ------------------ Basic PPT Generator (KEPT FOR COMPLETENESS, BUT ENHANCED IS RECOMMENDED) ------------------
def generate_ppt(message: str, df: pd.DataFrame, include_charts: bool = False):
    actions = extract_actions(message)
    descriptive_action = select_descriptive_action(actions)
    additional_actions = [a for a in actions if a != descriptive_action]

    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "🚀 Supply Sense AI"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = descriptive_action

    # Insights & Recommendations
    numeric_cols = df.select_dtypes(include='number').columns.tolist()
    categorical_cols = df.select_dtypes(exclude='number').columns.tolist()

    insights, recs = [], []

    for col in numeric_cols[:5]:
        insights.append(f"📊 '{col}': min={df[col].min()}, max={df[col].max()}, avg={df[col].mean():.2f}")
        recs.append(f"✅ Optimize '{col}' based on trends.")

    for col in categorical_cols[:3]:
        top_items = ', '.join(df[col].value_counts().head(3).index.astype(str))
        insights.append(f"🌟 Top '{col}': {top_items}")
        recs.append(f"✅ Focus on top '{col}' items: {top_items}.")

    if additional_actions:
        insights.extend([f"🔹 {act}" for act in additional_actions[:5]])

    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "📊 Insights & Recommendations"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.add_paragraph(); p.text = '\n'.join(insights); p.font.size = Pt(12)
    p = tf.add_paragraph(); p.text = '\n'.join(recs); p.font.size = Pt(12)

    # Chart Slides
    if include_charts and numeric_cols:
        for col in numeric_cols:
            plt.figure(figsize=(6,4))
            plt.bar(df.index.astype(str), df[col], color="#4CAF50")
            plt.title(f"{col} Analysis", fontsize=14)
            plt.xticks(rotation=45, ha='right')
            plt.tight_layout()
            chart_path = f"chart_{uuid.uuid4().hex[:8]}.png"
            plt.savefig(chart_path)
            plt.close()
            slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(chart_path, Inches(1), Inches(1.2), width=Inches(7))
            title_shape = slide.shapes.title if slide.shapes.title else slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5)
            )
            title_shape.text = f"📈 {col} Chart"
            if os.path.exists(chart_path):
                os.remove(chart_path)

    # Data Preview (safe)
    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "🗂 Data Preview"

    if df is not None and not df.empty and len(df.columns) > 0:
        display_cols = df.columns[:6]
        rows, cols = min(10, len(df)) + 1, len(display_cols)
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8)).table
        for col_idx, col_name in enumerate(display_cols):
            table.cell(0, col_idx).text = str(col_name)
        for row_idx, row in df.head(10).iterrows():
            for col_idx, col in enumerate(display_cols):
                table.cell(row_idx + 1, col_idx).text = str(row[col])
    else:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = "⚠️ No data available for preview."

    # Highlights
    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "📖 Key Highlights"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.word_wrap = True
    highlights = []
    for col in df.select_dtypes(include='number').columns[:5]:
        highlights.append(f"📌 Max {col}: {df[col].max()}, Min {col}: {df[col].min()}")
    for col in df.select_dtypes(exclude='number').columns[:3]:
        top_items = ', '.join(df[col].value_counts().head(3).index.astype(str))
        highlights.append(f"📌 Top '{col}': {top_items}")
    for h in highlights:
        p = tf.add_paragraph(); p.text = h; p.font.size = Pt(12)

    # Thank You
    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "🙏 Thank You"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "This presentation was auto-generated dynamically."

    # Save
    output_dir = "generated_files"
    os.makedirs(output_dir, exist_ok=True)
    filename = generate_unique_blob_name(prefix="PPT", actions_list=actions, extension="pptx")
    ppt_path = os.path.join(output_dir, filename)
    prs.save(ppt_path)
    return ppt_path

# ------------------ Corrected generate_excel (Only the enhanced one remains) ------------------
def generate_excel(df: pd.DataFrame, question: str, include_charts: bool = False, filename: str = None) -> str:
    """
    Generates an Excel file from a DataFrame.

    Uses the provided filename if available, otherwise creates a new one.
    """
    output_dir = "generated_files"
    os.makedirs(output_dir, exist_ok=True)
    
    # Use the provided filename or create a fallback
    #output_filename = filename if filename else f"report_{datetime.now().strftime('%Y%m%d%H%M%S')}.xlsx"
    if filename:
        base_name = os.path.splitext(filename)[0]
        output_filename = f"{base_name}.xlsx"
    else:
        output_filename = f"report_{datetime.now().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex[:8]}.xlsx"
    excel_path = os.path.join(output_dir, output_filename)
    
    # Save the DataFrame to the Excel file
    df.to_excel(excel_path, index=False)
    
    return excel_path

def generate_word(df: pd.DataFrame, question: str, include_charts: bool = False, filename: str = None) -> str:
    """
    Generates a Word document (.docx) from a DataFrame, including insights.

    Uses the provided filename if available, otherwise creates a new one.
    """
    doc = Document()
    doc.add_heading(question, 0)

    # Add data table (safe)
    if df is not None and not df.empty and len(df.columns) > 0:
        table = doc.add_table(rows=1, cols=len(df.columns))
        table.style = 'Table Grid'  # Add a simple style
        hdr_cells = table.rows[0].cells
        for i, col in enumerate(df.columns):
            hdr_cells[i].text = str(col)
        
        # Add data rows (up to 10)
        for _, row in df.head(10).iterrows():
            row_cells = table.add_row().cells
            for i, col in enumerate(df.columns):
                row_cells[i].text = str(row[col])
    else:
        doc.add_paragraph("No data available to display.")

    # Add GPT-generated insights and recommendations
    insights, recs = generate_insights(df)
    doc.add_heading('Insights', level=1)
    doc.add_paragraph(insights)
    doc.add_heading('Recommendations', level=1)
    doc.add_paragraph(recs)

    # --- Save File ---
    output_dir = "generated_files"
    os.makedirs(output_dir, exist_ok=True)
    
    # Use the provided filename or create a fallback
    #output_filename = filename if filename else f"report_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx"
    # Normalize filename to prevent double extensions
    
    if filename:
        # Strip extension safely to prevent double extensions
        base_name = os.path.splitext(filename)[0]
        output_filename = f"{base_name}.docx"
    else:
        # Fallback if no filename is provided
        output_filename = f"report_{datetime.now().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex[:8]}.docx"
    word_path = os.path.join(output_dir, output_filename)
    doc.save(word_path)
    
    return word_path

def generate_direct_response(question: str, df: pd.DataFrame) -> str:
    prompt = f"Answer the following question based on the data provided:\n\nQuestion: {question}\n\nData:\n{df.head(10).to_string(index=False)}"
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a helpful assistant providing direct answers based on data."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.2,
    )
    return response.choices[0].message.content.strip()

# ========================== ENHANCED, PROFESSIONAL PPT GENERATOR ==========================
def generate_ppt_enhanced(question: str, df: pd.DataFrame, include_charts: bool = True, filename: str = None):
    """
    Generate professional PowerPoint presentations with GPT-based summarized insights,
    auto-bulleted executive summaries, readable charts, and safe layouts.
    """

    from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR


    
    # --- Initialize Presentation ---
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def extract_bold_segments(text):
            """Return list of (segment, is_bold)."""
            pattern = r"\*\*(.*?)\*\*"
            result = []
            last_end = 0
            for match in re.finditer(pattern, text):
                normal = text[last_end:match.start()]
                if normal.strip():
                    result.append((normal.strip(), False))

                bold = match.group(1).strip()
                result.append((bold, True))

                last_end = match.end()

            if last_end < len(text):
                tail = text[last_end:].strip()
                if tail:
                    result.append((tail, False))

            return result


    # --- Helper: GPT short summarizer for title ---
    def summarize_title(text):
        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You create short, professional titles for PowerPoint reports."},
                    {"role": "user", "content": f"Summarize this in max 8 words for a PowerPoint title: {text}"}
                ],
                temperature=0.3,
            )
            return resp.choices[0].message.content.strip()
        except Exception:
            return text[:60]

    # --- Helper: GPT compression to bullets ---
    def compress_text_to_bullets(text, max_bullets=6):
        if not text or len(text.strip()) == 0:
            return "No insights available."
        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You summarize analytical text into clear, concise business bullet points for PowerPoint slides."},
                    {"role": "user", "content": f"Convert the following into {max_bullets} concise bullet points:\n{text}"}
                ],
                temperature=0.3,
            )
            # Remove Markdown list characters and return
            output = resp.choices[0].message.content.strip()
            output = re.sub(r"^\s*[-*]\s*", "", output, flags=re.MULTILINE)
            return output
        except Exception:
            return text[:1000]

    # --- Helper: Add footer brand line ---
    def add_footer(slide):
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        tf = footer.text_frame
        tf.text = "Supply Sense AI – Auto-generated Report"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

    # --- Helper: Add wrapped text safely ---
    def add_wrapped_text(slide, title, content, max_chars=800, font_size=18, top=1.0, max_lines=12):
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(28)
        title_box.text_frame.paragraphs[0].font.bold = True

        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(top), slide_width - Inches(2), slide_height - Inches(top + 1)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP

        # Ensure content is a list of lines/bullets for proper formatting
        if isinstance(content, str):
            lines = [line.strip() for line in content.split('\n') if line.strip()]
        else:
            lines = content

        adjusted_font = max(14, font_size - len(" ".join(lines)) // 200)

        # Content - format as bullet points
        tf.clear()
        
        # Add title text placeholder
        p = tf.add_paragraph()
        p.text = "Key Takeaways:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        # Add bullet points
        for i, line in enumerate(lines):
            if i >= max_lines:
                p = tf.add_paragraph()
                p.text = "..."
                p.font.size = Pt(adjusted_font)
                break
            p = tf.add_paragraph()
            segments = extract_bold_segments(line)
            p = tf.add_paragraph()
            p.level = 1

            for seg, is_bold in segments:
                run = p.add_run()
                run.text = seg + " "
                run.font.bold = is_bold
                run.font.size = Pt(adjusted_font)

            p.level = 1
            p.font.name = "Calibri"

        add_footer(slide)
        return textbox

    # --- Helper: Split long text into multiple slides ---
    def add_multislide_text(prs, title, content, max_bullets_per_slide=6):
        if isinstance(content, str):
            all_bullets = [line.strip() for line in content.split('\n') if line.strip()]
        else:
            all_bullets = content
            
        chunks = [all_bullets[i:i + max_bullets_per_slide] for i in range(0, len(all_bullets), max_bullets_per_slide)]
        
        for idx, chunk in enumerate(chunks, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_text = f"{title} (Part {idx})" if len(chunks) > 1 else title
            add_wrapped_text(slide, title_text, chunk, font_size=16, top=1.5, max_lines=max_bullets_per_slide)

    # --- 1. Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    title_box.text_frame.paragraphs[0].text = "🚀 Supply Sense AI"
    title_box.text_frame.paragraphs[0].font.size = Pt(40)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), slide_width - Inches(2), Inches(1.5))
    subtitle_box.text_frame.paragraphs[0].text = summarize_title(question)
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_footer(slide)

    # --- 2. Generate GPT Insights ---
    insights, recs = generate_insights(df)
    insight_bullets = compress_text_to_bullets(insights, max_bullets=6)
    rec_bullets = compress_text_to_bullets(recs, max_bullets=5)

    # --- 3. Executive Summary (Insights) ---
    add_multislide_text(prs, "📘 Executive Summary", insight_bullets)

    # --- 4. Recommendations ---
    add_multislide_text(prs, "💡 Strategic Recommendations", rec_bullets)

    # --- 5. Charts (max 3) ---
    if include_charts and df is not None and not df.empty:
        numeric_cols = df.select_dtypes(include="number").columns[:3]
        for col in numeric_cols:
            if len(df) > 0:
                plt.figure(figsize=(5, 3))
                # Use value counts for categorical or top N for numeric if time-series is unclear
                if df[col].nunique() < 20 and df[col].dtype != "number":
                    chart_data = df[col].value_counts().head(10)
                    plt.bar(chart_data.index.astype(str), chart_data.values, color="#007ACC")
                    plt.title(f"Top Categories of {col}")
                    
                else:
                    # Time-series like plot for numeric data (first 10 rows)
                    plt.plot(df.index.astype(str)[:10], df[col].head(10), marker='o', color="#4CAF50")
                    plt.title(f"Trend Snapshot: {col}")
                    
                plt.xticks(rotation=45, ha="right", fontsize=8)
                plt.yticks(fontsize=8)
                plt.tight_layout()
                path = f"chart_{uuid.uuid4().hex[:8]}.png"
                plt.savefig(path)
                plt.close()

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
                title_box.text_frame.paragraphs[0].text = f"📈 Data Visual: {col}"
                title_box.text_frame.paragraphs[0].font.size = Pt(22)
                title_box.text_frame.paragraphs[0].font.bold = True
                
                # Add image centered
                left = (slide_width - Inches(7)) / 2
                top = Inches(1.5)
                slide.shapes.add_picture(path, left, top, Inches(7), Inches(4))
                
                add_footer(slide)
                os.remove(path)

    # --- 6. Data Snapshot (Safe Table) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "🗂 Data Snapshot (Top 5 Rows)"

    if df is not None and not df.empty and len(df.columns) > 0:
        display_cols = df.columns[:6]
        rows = min(5, len(df)) + 1
        table = slide.shapes.add_table(rows, len(display_cols), Inches(0.5), Inches(1.2), Inches(9), Inches(3.5)).table
        
        # Style Header
        for i, col in enumerate(display_cols):
            cell = table.cell(0, i)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 102, 204) # Blue
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            
        # Add Data
        for r, (_, row) in enumerate(df.head(5).iterrows(), 1):
            for c, col in enumerate(display_cols):
                table.cell(r, c).text = str(row[col])[:30]
                table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(12)
    else:
        msg_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
        tf = msg_box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "⚠️ No data available to display.\nEnsure valid query results before generating a presentation."
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(200, 0, 0)

    add_footer(slide)

    # --- 7. Thank You Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ty_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(6), Inches(1.5))
    ty_box.text_frame.paragraphs[0].text = "🙏 Thank You"
    ty_box.text_frame.paragraphs[0].font.size = Pt(40)
    ty_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
 
    
    
    add_footer(slide)

    # --- 8. Save File ---
   
    # --- 8. Save File ---
    os.makedirs("generated_files", exist_ok=True)

    # Use provided filename or generate fallback
    # FIX: Use the filename directly if provided, otherwise generate a full, new filename.
    #output_filename = filename if filename else f"report_{datetime.now().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex[:8]}.pptx"

    # Normalize filename to prevent double extensions
    # Normalize filename to prevent double extensions
    
    if filename:
        base_name = os.path.splitext(filename)[0]
        output_filename = f"{base_name}.pptx"
    else:
        output_filename = f"report_{datetime.now().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex[:8]}.pptx"
        
    ppt_path = os.path.join("generated_files", output_filename)
    prs.save(ppt_path)
    return ppt_path
